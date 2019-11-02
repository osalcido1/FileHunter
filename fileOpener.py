from docx import Document
from docx.enum.text import WD_COLOR_INDEX
from pptx import Presentation
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.enum.dml import MSO_THEME_COLOR_INDEX
from pptx.dml.color import RGBColor
import openpyxl
from openpyxl import Workbook
from openpyxl import load_workbook
from openpyxl.styles import PatternFill, Border, Side, Alignment, Protection, Font
import os


def open_file(filename, target):
    if '.docx' in filename:
       doc = Document(filename)
       for paragraph in doc.paragraphs:
           if target in paragraph.text:
               for run in paragraph.runs:
                   if target in run.text:
                       if run.text == target:
                           run.font.highlight_color = WD_COLOR_INDEX.YELLOW
                           continue
                       x = run.text.split(target)
                       run.clear()
                       run.add_text(x[0])
                       temp = paragraph.add_run("", run.style)
                       temp.add_text(target)
                       temp.font.highlight_color = WD_COLOR_INDEX.YELLOW
                       temp2 = paragraph.add_run(x[1],run.style)


       doc.save(filename)
    if '.pptx' in filename:
       prs = Presentation (filename)
       text_runs= []
       for slides in prs.slides:
           for shape in slides.shapes:
               if not shape.has_text_frame:
                   continue
               text_frame = shape.text_frame
               for paragraph in text_frame.paragraphs:
                   for run in paragraph.runs:
                       if target in run.text:
                           if run.text == target:
                               run.font.fill.solid()
                               d = RGBColor(0xff, 0xff, 0x00)
                               run.font.fill.fore_color.rgb = d
                               continue
                           x = run.text.split(target);
                           run.text = x[0]
                           temp = paragraph.add_run()
                           temp.text = target
                           temp.font.bold = run.font.bold
                           temp.font.italic = run.font.italic
                           temp.font.language_id = run.font.language_id
                           temp.font.name = run.font.name
                           temp.font.size = run.font.size
                           temp.font.underline =temp.font.underline
                           temp.font.fill.solid()
                           d = RGBColor(0xff, 0xff, 0x00)
                           temp.font.fill.fore_color.rgb = d
                           temp2 = paragraph.add_run()
                           temp2.font.bold = run.font.bold
                           temp2.font.italic = run.font.italic
                           temp2.font.language_id = run.font.language_id
                           temp2.font.name = run.font.name
                           temp2.font.size = run.font.size
                           temp2.font.underline = temp.font.underline
                           temp2.text = x[1]

       prs.save(filename)
    if '.xlsx' in filename:
        wb = load_workbook(filename)
        yellowFill = PatternFill(start_color='FFFFFF00', end_color='FFFFFF00', fill_type='solid')
        for ws in wb.worksheets:
            for row in ws.iter_rows():
                for cell in row:
                    if cell.value == None:
                        continue
                    if type(cell.value) != str :
                        temp = str(cell.value)
                        if target in temp:
                            cell.fill = yellowFill
                        continue
                    if target in cell.value:
                        cell.fill = yellowFill
        wb.save(filename)

    os.startfile(filename)

open_file(r'C:\Users\Owen\PycharmProjects\FileHunter\venv\Include\test.xlsx', 'Dog')